# -*- coding: utf-8 -*-
"""pptx_to_text.py — convert 1 file .pptx thành text có cấu trúc cho lesson-prep.

Mỗi slide 1 khối, giữ [TITLE]/[BODY]/[TABLE]/[NOTES]. Ghi cạnh file gốc:
{file}.pptx.txt. In 1 dòng: PPTX <out> <slides>.

Chạy:  python pptx_to_text.py <file.pptx>
Cần: python-pptx.
"""
import sys, os

try:
    from pptx import Presentation
except ImportError:
    print("ERROR NOPPTX"); sys.exit(3)


def slide_to_text(slide, idx):
    lines = ["===== SLIDE %d =====" % idx]
    title, title_id = None, None
    try:
        t = slide.shapes.title
        if t is not None and (t.text or "").strip():
            title = t.text.strip(); title_id = t.shape_id
    except Exception:
        title = None
    if title:
        lines.append("[TITLE] " + title)
    for shape in slide.shapes:
        if title_id is not None and getattr(shape, "shape_id", None) == title_id:
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join((c.text or "").strip() for c in row.cells))
            if rows:
                lines.append("[TABLE] " + rows[0])
                lines.extend("        " + r for r in rows[1:])
        elif getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = (shape.text_frame.text or "").strip()
            if txt:
                lines.append("[BODY] " + txt.replace("\n", " / "))
    try:
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append("[NOTES] " + notes)
    except Exception:
        pass
    return "\n".join(lines)


def main(argv):
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")
        except (AttributeError, ValueError):
            pass
    if len(argv) != 2:
        print("Usage: python pptx_to_text.py <file.pptx>", file=sys.stderr); return 2
    path = argv[1]
    if not os.path.exists(path):
        print("ERROR NOFILE " + path); return 2
    prs = Presentation(path)
    blocks = [slide_to_text(s, i) for i, s in enumerate(prs.slides, start=1)]
    out = path + ".txt"
    with open(out, "w", encoding="utf-8") as f:
        f.write("\n\n".join(blocks))
    print("PPTX %s %d" % (out, len(blocks)))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
