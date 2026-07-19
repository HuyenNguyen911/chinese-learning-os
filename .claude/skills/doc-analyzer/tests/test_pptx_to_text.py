# -*- coding: utf-8 -*-
import os, sys, subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

HERE = Path(__file__).resolve().parent
SCRIPT = HERE.parent / "pptx_to_text.py"
PY = sys.executable
sys.path.insert(0, str(HERE.parent))
from pptx_to_text import _iter_shapes  # noqa: E402

def _make_pptx(path):
    prs = Presentation()
    # Slide 1: title + body
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "生词"
    s.placeholders[1].text = "尴尬 gāngà bối rối"
    s.notes_slide.notes_text_frame.text = "ghi chú của cô"
    # Slide 2: table
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = s2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    tbl.cell(0, 0).text = "词"; tbl.cell(0, 1).text = "拼音"
    tbl.cell(1, 0).text = "尴尬"; tbl.cell(1, 1).text = "gāngà"
    prs.save(str(path))

def test_pptx_to_text(tmp_path):
    pptx = tmp_path / "bai.pptx"
    _make_pptx(pptx)
    r = subprocess.run([PY, str(SCRIPT), str(pptx)], capture_output=True,
                       text=True, encoding="utf-8")
    assert r.returncode == 0, r.stderr
    assert r.stdout.startswith("PPTX ")
    out = pptx.with_suffix(".pptx.txt")
    assert out.exists()
    text = out.read_text(encoding="utf-8")
    assert "===== SLIDE 1 =====" in text
    assert "[TITLE] 生词" in text
    assert "尴尬" in text
    assert "[NOTES] ghi chú của cô" in text
    assert "===== SLIDE 2 =====" in text
    assert "[TABLE]" in text

def test_pptx_to_text_nofile(tmp_path):
    r = subprocess.run([PY, str(SCRIPT), str(tmp_path / "missing.pptx")],
                       capture_output=True, text=True, encoding="utf-8")
    assert "ERROR NOFILE" in r.stdout


class _Leaf:
    def __init__(self, n):
        self.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        self.n = n


class _Group:
    def __init__(self, kids):
        self.shape_type = MSO_SHAPE_TYPE.GROUP
        self.shapes = kids


def test_iter_shapes_recurses_nested_groups():
    # Group lồng group — phải "làm phẳng" theo đúng thứ tự, không bỏ sót shape trong group.
    tree = [_Leaf(1), _Group([_Leaf(2), _Group([_Leaf(3)]), _Leaf(4)]), _Leaf(5)]
    assert [s.n for s in _iter_shapes(tree)] == [1, 2, 3, 4, 5]
